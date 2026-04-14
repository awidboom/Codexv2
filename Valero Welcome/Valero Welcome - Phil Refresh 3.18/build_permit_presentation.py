from __future__ import annotations

import io
import json
import math
import re
from collections import Counter
from pathlib import Path
from typing import Iterable

import fitz
from PIL import Image, ImageDraw
from openpyxl import load_workbook
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parent
OUTPUT_PPTX = ROOT / "permit_significant_changes_presentation.pptx"
FALLBACK_PPTX = ROOT / "permit_significant_changes_presentation_updated.pptx"
ASSET_DIR = ROOT / "presentation_assets"
EXCERPT_DIR = ASSET_DIR / "excerpts"
OLD_PDF = ROOT / "09100062-103 Permit.pdf"
NEW_PDF = ROOT / "Preliminary Draft Permit.pdf"
WORKBOOK = ROOT / "permit_comparison_analysis.xlsx"
WORKBOOK_FALLBACK = ROOT / "permit_comparison_analysis_updated.xlsx"
REVIEW_JS = ROOT / "permit_review_data.js"

BG = RGBColor(245, 250, 253)
PANEL = RGBColor(255, 255, 255)
NAVY = RGBColor(24, 40, 54)
GOLD = RGBColor(4, 134, 204)
SLATE = RGBColor(128, 149, 158)
LIGHT = RGBColor(195, 210, 219)
RED = RGBColor(254, 115, 73)
GREEN = RGBColor(96, 183, 230)
FONT_MAJOR = "Arial"
FONT_MINOR = "Arial"


def load_review_rows() -> list[dict]:
    text = REVIEW_JS.read_text(encoding="utf-8")
    prefix = "window.PERMIT_REVIEW_DATA = "
    return json.loads(text[len(prefix):].rstrip().rstrip(";"))["rows"]


def load_workbook_rows() -> tuple[list[dict], Counter, Counter]:
    workbook_path = WORKBOOK_FALLBACK if WORKBOOK_FALLBACK.exists() and WORKBOOK_FALLBACK.stat().st_mtime >= WORKBOOK.stat().st_mtime else WORKBOOK
    wb = load_workbook(workbook_path, read_only=True, data_only=True)
    ws = wb["Detailed Comparison"]
    rows_iter = ws.iter_rows(values_only=True)
    headers = list(next(rows_iter))
    rows = [dict(zip(headers, row)) for row in rows_iter]
    change_counts = Counter(r["Change Type"] for r in rows)
    conclusion_counts = Counter(r["Consultant Conclusion"] for r in rows)
    return rows, change_counts, conclusion_counts


def normalize_space(text: str) -> str:
    return re.sub(r"\s+", " ", (text or "")).strip()


def summarize(text: str, limit: int) -> str:
    text = normalize_space(text)
    if len(text) <= limit:
        return text
    cut = text[: limit - 1].rsplit(" ", 1)[0]
    return cut + "..."


def slide_bullets(row: dict, raw_row: dict) -> list[tuple[str, str]]:
    return [
        ("Change", summarize(row["coreChange"], 110)),
        ("Why it matters", summarize(raw_row["Regulatory or Compliance Impact"], 110)),
        ("Ask", summarize(row["recommendedFollowUp"], 90)),
    ]


def chunk_lines(items: Iterable[str], max_chars: int) -> list[str]:
    lines: list[str] = []
    current = ""
    for item in items:
        item = item.strip()
        if not item:
            continue
        proposal = item if not current else f"{current}; {item}"
        if len(proposal) <= max_chars:
            current = proposal
        else:
            if current:
                lines.append(current)
            current = item
    if current:
        lines.append(current)
    return lines


def parse_row_lookup() -> dict[str, dict]:
    workbook_rows, _, _ = load_workbook_rows()
    return {row["Comparison ID"]: row for row in workbook_rows}


def search_text_rect(page: fitz.Page, text: str) -> fitz.Rect | None:
    clean = normalize_space(text)
    if not clean:
        return None
    words = clean.split()
    snippets: list[str] = []
    for size in (10, 8, 6):
        if len(words) >= size:
            positions = [0, max(0, len(words) // 3), max(0, len(words) // 2)]
            for pos in positions:
                snippet = " ".join(words[pos : pos + size])
                if len(snippet) > 20:
                    snippets.append(snippet)
    seen: set[str] = set()
    for snippet in snippets:
        if snippet in seen:
            continue
        seen.add(snippet)
        rects = page.search_for(snippet)
        if rects:
            return rects[0]
    return None


def union_rects(rects: list[fitz.Rect]) -> fitz.Rect:
    current = fitz.Rect(rects[0])
    for rect in rects[1:]:
        current |= rect
    return current


def excerpt_clip(
    page: fitz.Page,
    row_text: str,
    normalized_rects: list[dict],
) -> tuple[fitz.Rect, list[fitz.Rect]]:
    page_rect = page.rect
    rects: list[fitz.Rect] = []
    for item in normalized_rects or []:
        rects.append(
            fitz.Rect(
                page_rect.x0 + item["x"] * page_rect.width,
                page_rect.y0 + item["y"] * page_rect.height,
                page_rect.x0 + (item["x"] + item["w"]) * page_rect.width,
                page_rect.y0 + (item["y"] + item["h"]) * page_rect.height,
            )
        )
    if not rects:
        found = search_text_rect(page, row_text)
        if found:
            rects = [found]
    if rects:
        focus = union_rects(rects)
        text_height = min(230, 55 + len(normalize_space(row_text)) * 0.12)
        y0 = max(page_rect.y0, focus.y0 - 18)
        y1 = min(page_rect.y1, max(focus.y1 + text_height, focus.y0 + 120))
        x0 = max(page_rect.x0, page_rect.x0 + page_rect.width * 0.055)
        x1 = min(page_rect.x1, page_rect.x0 + page_rect.width * 0.965)
        clip = fitz.Rect(x0, y0, x1, y1)
        return clip, rects
    fallback_top = page_rect.y0 + page_rect.height * 0.14
    fallback_bottom = min(page_rect.y1, fallback_top + page_rect.height * 0.30)
    clip = fitz.Rect(
        page_rect.x0 + page_rect.width * 0.055,
        fallback_top,
        page_rect.x0 + page_rect.width * 0.965,
        fallback_bottom,
    )
    return clip, []


def render_excerpt_image(
    pdf_path: Path,
    page_num: int | None,
    row_text: str,
    normalized_rects: list[dict],
    output_path: Path,
) -> bool:
    if not page_num:
        return False
    doc = fitz.open(pdf_path)
    try:
        page = doc[page_num - 1]
        clip, rects = excerpt_clip(page, row_text, normalized_rects)
        matrix = fitz.Matrix(2.6, 2.6)
        pix = page.get_pixmap(matrix=matrix, clip=clip, alpha=False)
        image = Image.open(io.BytesIO(pix.tobytes("png"))).convert("RGB")
        draw = ImageDraw.Draw(image)
        scale_x = image.width / clip.width
        scale_y = image.height / clip.height
        for rect in rects:
            x0 = max(0, (rect.x0 - clip.x0) * scale_x - 8)
            y0 = max(0, (rect.y0 - clip.y0) * scale_y - 6)
            x1 = min(image.width, (rect.x1 - clip.x0) * scale_x + 8)
            y1 = min(image.height, (rect.y1 - clip.y0) * scale_y + 6)
            draw.rounded_rectangle(
                [x0, y0, x1, y1],
                radius=8,
                outline=(254, 115, 73),
                width=5,
                fill=(223, 241, 251),
            )
        output_path.parent.mkdir(parents=True, exist_ok=True)
        image.save(output_path)
        return True
    finally:
        doc.close()


def set_background(slide) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG


def add_textbox(slide, left, top, width, height, text, *, font_size=18, bold=False,
                color=NAVY, align=PP_ALIGN.LEFT, name=FONT_MINOR, margins=(0.06, 0.04, 0.04, 0.02)):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(margins[0])
    tf.margin_right = Inches(margins[1])
    tf.margin_top = Inches(margins[2])
    tf.margin_bottom = Inches(margins[3])
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_panel(slide, left, top, width, height, *, fill_color=PANEL, line_color=LIGHT):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = line_color
    shape.line.width = Pt(1.2)
    return shape


def add_bullets(slide, left, top, width, height, items: list[str], *, font_size=18, color=SLATE):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.06)
    tf.margin_right = Inches(0.04)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    first = True
    for item in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        p.level = 0
        p.text = item
        p.font.name = FONT_MINOR
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.bullet = True
        p.space_after = Pt(4)
        first = False
    return box


def fit_picture(slide, image_path: Path, left, top, width, height):
    with Image.open(image_path) as img:
        img_ratio = img.width / img.height
    box_ratio = width / height
    if img_ratio >= box_ratio:
        pic_w = width
        pic_h = width / img_ratio
        pic_left = left
        pic_top = top + (height - pic_h) / 2
    else:
        pic_h = height
        pic_w = height * img_ratio
        pic_top = top
        pic_left = left + (width - pic_w) / 2
    slide.shapes.add_picture(str(image_path), pic_left, pic_top, pic_w, pic_h)


def add_metric(slide, left, top, width, height, label, value, accent):
    panel = add_panel(slide, left, top, width, height, fill_color=PANEL, line_color=LIGHT)
    panel.fill.fore_color.rgb = PANEL
    bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, left, top, Inches(0.12), height)
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent
    bar.line.fill.background()
    add_textbox(slide, left + Inches(0.18), top + Inches(0.08), width - Inches(0.24), Inches(0.25),
                label, font_size=13, bold=True, color=SLATE)
    add_textbox(slide, left + Inches(0.18), top + Inches(0.34), width - Inches(0.24), Inches(0.36),
                str(value), font_size=24, bold=True, color=NAVY)


def build_deck() -> Path:
    review_rows = load_review_rows()
    row_lookup = parse_row_lookup()
    workbook_rows, change_counts, conclusion_counts = load_workbook_rows()
    focus_rows = [r for r in review_rows if r.get("reviewPriority") == "Highlight"]
    if not focus_rows:
        focus_rows = [r for r in review_rows if r["consultantConclusion"] == "Significant changes"]
    focus_rows.sort(key=lambda r: (r.get("importanceRank") or 999, r["id"]))
    subject_counts = Counter(r["Subject Tag Primary"] for r in workbook_rows)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    title = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(title)
    band = title.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.18))
    band.fill.solid()
    band.fill.fore_color.rgb = GOLD
    band.line.fill.background()
    add_textbox(title, Inches(0.65), Inches(0.55), Inches(8.4), Inches(0.6),
                "Valero Welcome Permit Draft Review", font_size=26, bold=True, color=NAVY)
    add_textbox(title, Inches(0.67), Inches(1.2), Inches(8.0), Inches(0.42),
                "Discussion deck focused on unrequested or unexplained draft changes", font_size=15, color=SLATE)
    ribbon = title.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.65), Inches(1.85), Inches(2.45), Inches(0.42))
    ribbon.fill.solid()
    ribbon.fill.fore_color.rgb = GOLD
    ribbon.line.fill.background()
    add_textbox(title, Inches(0.82), Inches(1.93), Inches(2.1), Inches(0.2),
                "Overall outcome: Significant changes", font_size=12, bold=True, color=PANEL)
    add_panel(title, Inches(0.65), Inches(2.55), Inches(12.0), Inches(3.8))
    metrics = [
        ("Comparison rows", len(workbook_rows), NAVY),
        ("Highlight items", len(focus_rows), RED),
        ("High-risk rows", sum(1 for r in workbook_rows if r["Risk Level"] == "High"), GOLD),
        ("High-confidence rows", sum(1 for r in workbook_rows if r.get("Evidence Confidence") == "High"), GREEN),
    ]
    left = Inches(0.95)
    for label, value, accent in metrics:
        add_metric(title, left, Inches(2.95), Inches(2.7), Inches(0.95), label, value, accent)
        left += Inches(2.9)
    add_textbox(title, Inches(0.95), Inches(4.25), Inches(4.0), Inches(0.28),
                "Primary themes", font_size=15, bold=True, color=NAVY)
    add_bullets(
        title,
        Inches(0.95),
        Inches(4.55),
        Inches(11.0),
        Inches(1.5),
        [
            "Modeling framework shifts from PM10 trigger language to boundary-control and public-access obligations.",
            "New ambient-boundary inspection, patrol, and breach-reporting requirements create recurring compliance tasks.",
            "Draft permit introduces multiple changes not clearly identified in the major amendment application narrative.",
            "Manual QA confirmed workflow false positives such as COMG 2, so this deck follows consultant-validated rows rather than raw workflow summary tabs.",
        ],
        font_size=16,
    )

    summary = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(summary)
    band = summary.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.14))
    band.fill.solid()
    band.fill.fore_color.rgb = GOLD
    band.line.fill.background()
    add_textbox(summary, Inches(0.65), Inches(0.45), Inches(6.3), Inches(0.4),
                "Executive Summary", font_size=24, bold=True, color=NAVY)
    add_panel(summary, Inches(0.65), Inches(1.0), Inches(4.1), Inches(5.75))
    add_textbox(summary, Inches(0.9), Inches(1.22), Inches(3.4), Inches(0.3),
                "Counts by change type", font_size=15, bold=True, color=NAVY)
    add_bullets(
        summary,
        Inches(0.9),
        Inches(1.58),
        Inches(3.4),
        Inches(1.55),
        [f"{k}: {change_counts[k]}" for k in ("Revision", "Addition", "Deletion", "Possible relocation or renumbering") if k in change_counts],
        font_size=16,
    )
    add_textbox(summary, Inches(0.9), Inches(3.0), Inches(3.4), Inches(0.3),
                "Evidence confidence", font_size=15, bold=True, color=NAVY)
    add_bullets(
        summary,
        Inches(0.9),
        Inches(3.32),
        Inches(3.4),
        Inches(0.7),
        [
            f"High confidence: {sum(1 for r in workbook_rows if r.get('Evidence Confidence') == 'High')}",
            f"Medium confidence: {sum(1 for r in workbook_rows if r.get('Evidence Confidence') == 'Medium')}",
        ],
        font_size=15,
    )
    add_textbox(summary, Inches(0.9), Inches(4.05), Inches(3.4), Inches(0.3),
                "Top subject areas", font_size=15, bold=True, color=NAVY)
    top_subjects = [f"{tag}: {count}" for tag, count in subject_counts.most_common(6)]
    add_bullets(summary, Inches(0.9), Inches(4.39), Inches(3.4), Inches(1.75), top_subjects[:4], font_size=15)
    add_panel(summary, Inches(4.95), Inches(1.0), Inches(7.7), Inches(5.75))
    add_textbox(summary, Inches(5.2), Inches(1.22), Inches(5.5), Inches(0.3),
                "Slides included in this deck", font_size=15, bold=True, color=NAVY)
    slide_items = [
        "Permit-wide and modeling changes replacing the existing PM10 framework",
        "Ambient boundary, fencing, breach-reporting, patrol, and camera obligations",
        "Unrequested STRU 11 language changes, including throughput and AOS applicability issues",
        "CAM and testing revisions not clearly requested in the application package",
    ]
    add_bullets(summary, Inches(5.2), Inches(1.58), Inches(6.9), Inches(2.4), slide_items, font_size=16)
    add_textbox(summary, Inches(5.2), Inches(4.35), Inches(6.9), Inches(0.8),
                "Use this deck for discussion. Use the spreadsheet for full analysis and line-item support.",
                font_size=17, color=SLATE)

    table = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(table)
    band = table.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.14))
    band.fill.solid()
    band.fill.fore_color.rgb = GOLD
    band.line.fill.background()
    add_textbox(table, Inches(0.65), Inches(0.45), Inches(8.0), Inches(0.4),
                "Significant Change Index", font_size=24, bold=True, color=NAVY)
    add_panel(table, Inches(0.65), Inches(1.0), Inches(12.0), Inches(5.9))
    col_x = [0.9, 1.7, 4.6, 8.7, 11.1]
    headers = ["ID", "Subject", "Core change", "Risk", "Follow-up"]
    widths = [0.6, 2.6, 4.0, 0.8, 1.2]
    for x, head, width in zip(col_x, headers, widths):
        add_textbox(table, Inches(x), Inches(1.2), Inches(width), Inches(0.25), head, font_size=12, bold=True, color=NAVY)
    y = 1.55
    for row in focus_rows[:10]:
        add_textbox(table, Inches(col_x[0]), Inches(y), Inches(widths[0]), Inches(0.4), row["id"], font_size=11, bold=True, color=GOLD)
        add_textbox(table, Inches(col_x[1]), Inches(y), Inches(widths[1]), Inches(0.52), row["subjectPrimary"], font_size=11, color=SLATE)
        add_textbox(table, Inches(col_x[2]), Inches(y), Inches(widths[2]), Inches(0.62), summarize(row["coreChange"], 105), font_size=11, color=SLATE)
        add_textbox(table, Inches(col_x[3]), Inches(y), Inches(widths[3]), Inches(0.3), row["riskLevel"], font_size=11, bold=True, color=RED if row["riskLevel"] == "High" else GOLD)
        add_textbox(table, Inches(col_x[4]), Inches(y), Inches(widths[4]), Inches(0.55), summarize(row["recommendedFollowUp"], 40), font_size=11, color=SLATE)
        line = table.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.92), Inches(y + 0.47), Inches(11.35), Pt(1))
        line.fill.solid()
        line.fill.fore_color.rgb = LIGHT
        line.line.fill.background()
        y += 0.52

    for row in focus_rows:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        set_background(slide)
        band = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.14))
        band.fill.solid()
        band.fill.fore_color.rgb = GOLD
        band.line.fill.background()
        raw_row = row_lookup[row["id"]]
        add_textbox(slide, Inches(0.55), Inches(0.28), Inches(9.2), Inches(0.34),
                    f"{row['id']} | {row['subjectPrimary']}", font_size=22, bold=True, color=NAVY)
        meta_line = f"Old: {row['oldSection'] or 'No direct equivalent'}    New: {row['newSection'] or 'No direct equivalent'}    Risk: {row['riskLevel']}"
        add_textbox(slide, Inches(0.58), Inches(0.64), Inches(7.5), Inches(0.22), meta_line, font_size=11, color=SLATE)
        add_textbox(slide, Inches(0.58), Inches(0.82), Inches(7.9), Inches(0.18),
                    f"Application basis: {row.get('requestedChangeStatus', 'N/A')} | {row.get('applicationSource', 'None')}",
                    font_size=10, color=SLATE)
        add_textbox(slide, Inches(8.25), Inches(0.82), Inches(2.0), Inches(0.18),
                    f"Evidence: {row.get('evidenceConfidence', 'N/A')}",
                    font_size=10, color=SLATE, align=PP_ALIGN.RIGHT)
        tag = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(10.55), Inches(0.30), Inches(2.15), Inches(0.38))
        tag.fill.solid()
        tag.fill.fore_color.rgb = RED if row["riskLevel"] == "High" else GOLD
        tag.line.fill.background()
        add_textbox(slide, Inches(10.75), Inches(0.37), Inches(1.8), Inches(0.16),
                    row.get("reviewPriority", row["consultantConclusion"]), font_size=11, bold=True, color=PANEL, align=PP_ALIGN.CENTER)

        add_panel(slide, Inches(0.55), Inches(1.02), Inches(6.05), Inches(4.49))
        add_panel(slide, Inches(6.73), Inches(1.02), Inches(6.05), Inches(4.49))
        add_textbox(slide, Inches(0.8), Inches(1.08), Inches(2.0), Inches(0.22),
                    f"Old permit | page {row['oldPage']}" if row.get("oldPage") else "Old permit", font_size=13, bold=True, color=NAVY)
        add_textbox(slide, Inches(6.98), Inches(1.08), Inches(2.0), Inches(0.22),
                    f"New permit | page {row['newPage']}" if row.get("newPage") else "New permit", font_size=13, bold=True, color=NAVY)

        old_img = EXCERPT_DIR / f"{row['id']}_old.png"
        new_img = EXCERPT_DIR / f"{row['id']}_new.png"
        has_old = render_excerpt_image(OLD_PDF, row.get("oldPage"), row.get("oldText", ""), row.get("oldHighlightRects", []), old_img)
        has_new = render_excerpt_image(NEW_PDF, row.get("newPage"), row.get("newText", ""), row.get("newHighlightRects", []), new_img)

        if has_old:
            fit_picture(slide, old_img, Inches(0.78), Inches(1.38), Inches(5.58), Inches(3.85))
        else:
            placeholder = add_panel(slide, Inches(0.93), Inches(1.78), Inches(5.28), Inches(2.6), fill_color=LIGHT)
            placeholder.line.color.rgb = LIGHT
            add_textbox(slide, Inches(1.15), Inches(2.62), Inches(4.9), Inches(0.5),
                        "No corresponding condition in the issued permit.", font_size=18, bold=True, color=SLATE, align=PP_ALIGN.CENTER)

        if has_new:
            fit_picture(slide, new_img, Inches(6.96), Inches(1.38), Inches(5.58), Inches(3.85))
        else:
            placeholder = add_panel(slide, Inches(7.11), Inches(1.78), Inches(5.28), Inches(2.6), fill_color=LIGHT)
            placeholder.line.color.rgb = LIGHT
            add_textbox(slide, Inches(7.33), Inches(2.62), Inches(4.9), Inches(0.5),
                        "No corresponding condition in the draft permit.", font_size=18, bold=True, color=SLATE, align=PP_ALIGN.CENTER)

        add_panel(slide, Inches(0.55), Inches(5.72), Inches(12.18), Inches(1.15))
        bullets = slide_bullets(row, raw_row)
        box_left = Inches(0.82)
        for label, text in bullets:
            add_panel(slide, box_left, Inches(5.9), Inches(3.8), Inches(0.72), fill_color=LIGHT, line_color=LIGHT)
            add_textbox(slide, box_left + Inches(0.12), Inches(6.0), Inches(0.95), Inches(0.16),
                        label, font_size=12, bold=True, color=NAVY)
            add_textbox(slide, box_left + Inches(0.12), Inches(6.18), Inches(3.5), Inches(0.28),
                        text, font_size=13, color=SLATE)
            box_left += Inches(4.02)

    closing = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(closing)
    band = closing.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.14))
    band.fill.solid()
    band.fill.fore_color.rgb = GOLD
    band.line.fill.background()
    add_textbox(closing, Inches(0.7), Inches(0.6), Inches(5.5), Inches(0.45),
                "Close-Out", font_size=24, bold=True, color=NAVY)
    add_panel(closing, Inches(0.7), Inches(1.2), Inches(12.0), Inches(5.6))
    add_textbox(closing, Inches(0.95), Inches(1.48), Inches(4.4), Inches(0.28),
                "Recommended client discussion points", font_size=15, bold=True, color=NAVY)
    add_bullets(
        closing,
        Inches(0.95),
        Inches(1.82),
        Inches(5.3),
        Inches(3.6),
        [
            "Whether the STRU 11 direct-vent AOS 2 is operationally necessary and fully supported by the application and modeling basis.",
            "Whether MPCA intends the deleted PM10 modeling trigger framework to be replaced by other enforceable documents or internal review practices.",
            "Whether new boundary security, patrol, camera, and breach-reporting obligations match current site controls without added capital or staffing.",
            "Whether revised throughput and CAM thresholds align with current operations and accepted October 30, 2025 test data.",
        ],
        font_size=17,
    )
    add_textbox(closing, Inches(6.8), Inches(1.48), Inches(4.4), Inches(0.28),
                "Files in the package", font_size=15, bold=True, color=NAVY)
    add_bullets(
        closing,
        Inches(6.8),
        Inches(1.82),
        Inches(4.9),
        Inches(2.4),
        [
            "permit_comparison_analysis.xlsx",
            "permit_comparison_summary.md / .docx",
            "permit_significant_changes_presentation.pptx",
            "permit_review_interface.html for QAQC",
        ],
        font_size=17,
    )
    add_textbox(closing, Inches(6.8), Inches(4.7), Inches(5.0), Inches(0.8),
                "Excerpt images are pulled from the mapped permit pages and flagged at the condition location when possible.",
                font_size=16, color=SLATE)

    ASSET_DIR.mkdir(parents=True, exist_ok=True)
    try:
        prs.save(OUTPUT_PPTX)
        return OUTPUT_PPTX
    except PermissionError:
        prs.save(FALLBACK_PPTX)
        return FALLBACK_PPTX


if __name__ == "__main__":
    output = build_deck()
    print(output)
