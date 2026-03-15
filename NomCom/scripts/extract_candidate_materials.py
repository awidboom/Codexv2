import argparse
import re
import zipfile
from pathlib import Path
from xml.etree import ElementTree as ET

import pandas as pd
from docx import Document
from pypdf import PdfReader
from pptx import Presentation


def _safe_filename(path: Path) -> str:
    return re.sub(r"[^A-Za-z0-9._ ()-]+", "_", path.name)


def extract_docx_to_text(path: Path) -> str:
    doc = Document(str(path))
    lines: list[str] = [f"FILE: {path.name}", ""]

    for paragraph in doc.paragraphs:
        text = (paragraph.text or "").strip()
        if text:
            lines.append(text)

    if doc.tables:
        lines.extend(["", "---- TABLES ----", ""])
        for table_index, table in enumerate(doc.tables, start=1):
            for row_index, row in enumerate(table.rows, start=1):
                for col_index, cell in enumerate(row.cells, start=1):
                    cell_paragraphs = []
                    for p in cell.paragraphs:
                        t = (p.text or "").strip()
                        if t:
                            cell_paragraphs.append(t)
                    if not cell_paragraphs:
                        continue
                    lines.append(f"[Table {table_index} Row {row_index} Col {col_index}]")
                    lines.extend(cell_paragraphs)
                    lines.append("")

    return "\n".join(lines).strip() + "\n"


def extract_dotx_to_text(path: Path) -> str:
    lines: list[str] = [f"FILE: {path.name}", ""]
    with zipfile.ZipFile(path, "r") as zf:
        try:
            xml_bytes = zf.read("word/document.xml")
        except KeyError:
            return "\n".join(lines).strip() + "\n"

    root = ET.fromstring(xml_bytes)
    texts: list[str] = []
    for node in root.iter():
        if node.tag.endswith("}t") and node.text:
            texts.append(node.text)
        elif node.tag.endswith("}tab"):
            texts.append("\t")
        elif node.tag.endswith("}br") or node.tag.endswith("}cr"):
            texts.append("\n")
        elif node.tag.endswith("}p"):
            texts.append("\n")

    raw = "".join(texts)
    normalized_lines = [" ".join(line.split()) for line in raw.splitlines()]
    normalized_lines = [line for line in normalized_lines if line]
    lines.extend(normalized_lines)
    return "\n".join(lines).strip() + "\n"


def extract_pptx_to_text(path: Path) -> str:
    prs = Presentation(str(path))
    lines: list[str] = [f"FILE: {path.name}", ""]
    for slide_index, slide in enumerate(prs.slides, start=1):
        lines.append(f"---- SLIDE {slide_index} ----")
        for shape in slide.shapes:
            if not hasattr(shape, "text"):
                continue
            text = " ".join((shape.text or "").split())
            if text:
                lines.append(text)
        lines.append("")
    return "\n".join(lines).strip() + "\n"


def extract_xlsx_to_text(path: Path) -> str:
    lines: list[str] = [f"FILE: {path.name}", ""]
    workbook = pd.ExcelFile(path)
    for sheet_name in workbook.sheet_names:
        df = workbook.parse(sheet_name=sheet_name, dtype=str, header=None)
        df = df.fillna("")
        lines.append(f"---- SHEET: {sheet_name} ----")
        for row in df.itertuples(index=False, name=None):
            row_cells = [str(cell).replace("\t", " ").replace("\r", " ").replace("\n", " ").strip() for cell in row]
            if any(row_cells):
                lines.append("\t".join(row_cells).rstrip())
        lines.append("")
    return "\n".join(lines).strip() + "\n"


def extract_pdf_to_text(path: Path) -> str:
    reader = PdfReader(str(path))
    lines: list[str] = [f"FILE: {path.name}", ""]
    for page_index, page in enumerate(reader.pages, start=1):
        text = (page.extract_text() or "").strip()
        if not text:
            continue
        lines.append(f"---- PAGE {page_index} ----")
        lines.append(text)
        lines.append("")
    return "\n".join(lines).strip() + "\n"


def write_text(out_path: Path, text: str) -> None:
    out_path.parent.mkdir(parents=True, exist_ok=True)
    out_path.write_text(text, encoding="utf-8")


def make_evalform_markdown(extracted_txt: str) -> str:
    lines_out: list[str] = []
    for raw in extracted_txt.splitlines():
        line = raw.strip()
        if not line:
            lines_out.append("")
            continue
        if line.startswith("FILE:"):
            continue
        if re.match(r"^(Background Summary|BACKGROUND SUMMARY)\b", line):
            lines_out.append("# " + line.title())
            continue
        if re.match(r"^Criteria\s*#?\s*[AP]\d\b", line, flags=re.IGNORECASE):
            lines_out.append("## " + line)
            continue
        lines_out.append(line)

    md = "\n".join(lines_out).strip() + "\n"
    if not md.startswith("#"):
        md = "# Extracted content: Evaluation Form\n\n" + md
    return md


def main() -> None:
    parser = argparse.ArgumentParser(description="Extract candidate Office/PDF materials into searchable text.")
    parser.add_argument("--input", required=True, help="Candidate folder (e.g., 'Candidate B').")
    parser.add_argument("--out", default=None, help="Output folder (default: '<input>/_extracted').")
    args = parser.parse_args()

    input_folder = Path(args.input).resolve()
    if not input_folder.exists():
        raise SystemExit(f"Input folder not found: {input_folder}")

    out_folder = Path(args.out).resolve() if args.out else (input_folder / "_extracted")
    out_folder.mkdir(parents=True, exist_ok=True)

    evalform_docx = None
    for candidate in input_folder.glob("**/*"):
        if candidate.is_file() and re.search(r"Principal VP Evaluation Form", candidate.name, flags=re.IGNORECASE):
            if candidate.suffix.lower() in {".docx", ".dotx"}:
                evalform_docx = candidate
                break

    for file_path in input_folder.glob("**/*"):
        if not file_path.is_file():
            continue
        suffix = file_path.suffix.lower()
        if suffix not in {".docx", ".dotx", ".pptx", ".xlsx", ".pdf"}:
            continue
        rel = file_path.relative_to(input_folder)
        out_path = out_folder / (str(rel) + ".txt")
        out_path.parent.mkdir(parents=True, exist_ok=True)

        if suffix == ".docx":
            text = extract_docx_to_text(file_path)
        elif suffix == ".dotx":
            text = extract_dotx_to_text(file_path)
        elif suffix == ".pptx":
            text = extract_pptx_to_text(file_path)
        elif suffix == ".xlsx":
            text = extract_xlsx_to_text(file_path)
        elif suffix == ".pdf":
            text = extract_pdf_to_text(file_path)
        else:
            continue

        write_text(out_path, text)

    if evalform_docx:
        eval_text_path = out_folder / (str(evalform_docx.relative_to(input_folder)) + ".txt")
        if eval_text_path.exists():
            md = make_evalform_markdown(eval_text_path.read_text(encoding="utf-8"))
            (input_folder / "_evalform_extracted.md").write_text(md, encoding="utf-8")


if __name__ == "__main__":
    main()
